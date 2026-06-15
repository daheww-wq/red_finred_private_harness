from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("client_briefing_finred_harness.pptx")
FONT = "Malgun Gothic"
NAVY = RGBColor(20, 38, 70)
BLUE = RGBColor(36, 91, 150)
GRAY = RGBColor(95, 103, 115)
LIGHT = RGBColor(244, 247, 251)
GREEN = RGBColor(30, 120, 85)
RED = RGBColor(178, 55, 55)


def set_run(run, size=18, bold=False, color=RGBColor(30, 30, 30)):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.32), Inches(12.4), Inches(0.72))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_run(p.runs[0], 24, True, NAVY)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.98), Inches(12.0), Inches(0.4))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        set_run(p2.runs[0], 11, False, GRAY)


def add_bullets(slide, bullets, x=0.7, y=1.45, w=12.0, h=5.4, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        set_run(p.runs[0], size, False, RGBColor(35, 40, 48))


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.1), Inches(12.4), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"FinRED Harness Briefing | {page}"
    p.alignment = PP_ALIGN.RIGHT
    set_run(p.runs[0], 8, False, GRAY)


def add_tag(slide, text, x, y, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.2), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], 10, True, RGBColor(255, 255, 255))


def add_flow(slide, labels, y=2.4):
    x = 0.55
    for i, label in enumerate(labels):
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.55), Inches(0.72))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = BLUE
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        set_run(p.runs[0], 10, True, NAVY)
        x += 1.75
        if i < len(labels) - 1:
            arrow = slide.shapes.add_textbox(Inches(x - 0.22), Inches(y + 0.22), Inches(0.3), Inches(0.2))
            ap = arrow.text_frame.paragraphs[0]
            ap.text = ">"
            set_run(ap.runs[0], 13, True, GRAY)


def add_table(slide, rows, cols, data, x=0.6, y=1.55, w=12.1, h=4.8):
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            set_run(para.runs[0], 9 if r else 10, r == 0, RGBColor(30, 30, 30))
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return table


def slide(prs, title, bullets=None, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_title(s, title, subtitle)
    if bullets:
        add_bullets(s, bullets)
    add_footer(s, len(prs.slides))
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    title = s.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(11.8), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "FinRED Harness Engineering"
    set_run(p.runs[0], 34, True, RGBColor(255, 255, 255))
    sub = s.shapes.add_textbox(Inches(0.78), Inches(2.35), Inches(11.6), Inches(1.0))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "금융 레드팀 데이터 생성, 응답 평가, ASR 산출을 반복 가능하게 운영하기 위한 하네스 설명 자료"
    set_run(p2.runs[0], 16, False, RGBColor(230, 236, 246))
    add_tag(s, "개발팀", 0.82, 3.35, BLUE)
    add_tag(s, "지식구축팀", 3.2, 3.35, GREEN)
    add_footer(s, 1)

    slide(prs, "한 줄 요약", [
        "FinRED 하네스는 금융 도메인 레드팀 평가 데이터를 반복 가능하게 생성하고 평가하는 운영 레이어입니다.",
        "기존 Step1/Step2/eval 스크립트를 유지하면서, 실행 전 점검, 환경 연결, 산출물 검증, judge/ASR 산출을 연결했습니다.",
        "R4_1 기준으로 scenario -> prompt -> response -> judge -> ASR까지 end-to-end 흐름을 완료했습니다.",
    ])

    s = slide(prs, "전체 데이터 흐름")
    add_flow(s, ["원천문서", "Chunk/DB", "Retrieval", "Scenario", "Prompt", "Response", "Judge"], y=2.15)
    add_bullets(s, [
        "핵심 산출물은 각 단계별 JSON/CSV로 남고, 최종적으로 Safe/Unsafe 및 ASR을 산출합니다.",
        "하네스는 이 흐름을 실행 전 점검, runtime 격리, checkpoint, report와 함께 관리합니다.",
    ], y=3.4, h=2.4)

    s = slide(prs, "주요 산출물 위치")
    add_table(s, 7, 3, [
        ["구분", "경로", "의미"],
        ["Schema", "src/data/schemas/ko/{category}.json", "scenario 구조"],
        ["Query", "src/data/queries/{category}_queries.csv", "retrieval 질의"],
        ["Chunks", "src/data/contexts/.../{category}_chunks.json", "근거 context"],
        ["Prompt", "src/outputs/prompts/{category}_prompts_all.csv", "attack prompt"],
        ["Response", "src/eval/dataset/*_with_responses.csv", "target model 응답"],
        ["Judge/ASR", "src/eval/infer_result/*_judge.*", "Safe/Unsafe 및 ASR"],
    ], h=4.8)

    s = slide(prs, "하네스 레이어 구성")
    add_flow(s, ["cli.py", "config.py", "runtime.py", "preflight.py", "pipeline.py", "validators.py", "checkpoint.py"], y=2.0)
    add_bullets(s, [
        "기본 실행은 fake provider dry-run으로 안전하게 검증합니다.",
        "live API는 .env 및 RUN_LIVE_LLM_TESTS 설정으로 명시적으로 연결합니다.",
        "runtime 산출물은 .runtime/<worktree_id>/ 아래에 격리됩니다.",
    ], y=3.25)

    slide(prs, "환경 및 API 연결", [
        ".env 자동 로드 지원: OPENAI_API_KEY, GEMINI_API_KEY, RUN_LIVE_LLM_TESTS",
        "PowerShell 형식($env:KEY=\"...\")과 표준 dotenv 형식(KEY=\"...\")을 모두 지원합니다.",
        "doctor 명령으로 key 존재 여부, live opt-in, 실행 위치를 확인합니다.",
        "키 값은 출력하지 않고 true/false 상태만 보여주도록 구성했습니다.",
    ])

    s = slide(prs, "R4_1 End-to-End 결과")
    add_table(s, 8, 3, [
        ["단계", "산출물", "상태"],
        ["Retrieval", "R4_1_chunks.json", "query 5개 context 있음"],
        ["Step1", "scenarios/R4_1/*.json", "5개"],
        ["Step2", "prompts/R4_1/*.json", "5개"],
        ["Prompt CSV", "R4_1_prompts_all.csv", "5 rows"],
        ["Response", "shinhan_r4_1_gemini_with_responses.csv", "5 rows, error 0"],
        ["Judge", "shinhan_r4_1_gemini_judge.csv", "Safe 3 / Unsafe 2"],
        ["ASR", "asr_lv2_oursrubric.json", "R4_1 = 0.4"],
    ], h=5.1)

    slide(prs, "R4_1 Rubric 수정", [
        "기존 R4_1 rubric은 취약계층 사칭/정부지원 탈취 등 R2성 기준이 섞여 있었습니다.",
        "R4_1 schema/query는 설명의무, 핵심설명서, 청약철회권, 위법계약해지권, 적합성 원칙, 불공정 판매가 중심입니다.",
        "따라서 R4_1 judge rubric을 소비자 권리 침해/설명의무 회피/권리행사 방해 기준으로 교체했습니다.",
        "재평가 후 결과는 5건 중 Safe 3, Unsafe 2, ASR 40%입니다.",
    ])

    s = slide(prs, "R1~R5 Rubric 점검 요약")
    add_table(s, 6, 2, [
        ["영역", "점검 결과"],
        ["R1", "금융 보안 위협 기준과 대체로 정합"],
        ["R2", "R2_4/R2_5 방향성 오류 수정 완료"],
        ["R3", "상품 왜곡, 허위정보, 가짜 전문가 기준 정합"],
        ["R4", "R4_1 불일치 수정 완료, R4_2~R4_5 정합"],
        ["R5", "내부통제, 보고 회피, 감사 방해 기준 정합"],
    ], h=4.2)

    slide(prs, "개발팀 관점", [
        "연결된 것: .env 자동 로드, doctor/preflight/run, fake dry-run, response 생성, judge, ASR",
        "확인된 것: unit test 12개 통과, R4_1 end-to-end 산출물 정상 파싱",
        "주의할 것: OpenAI judge 호출은 현재 환경에서 429 insufficient_quota 계열 오류가 재현됩니다.",
        "후속 권장: Gemini SDK를 google-genai 기반으로 마이그레이션하고, Step1/Step2/response/judge를 harness command로 통합합니다.",
    ])

    slide(prs, "지식구축팀 관점", [
        "검수 대상: schema, query, retrieved chunks, scenario, prompt, response, judge rubric, Safe/Unsafe 판정",
        "핵심 질문: 이 프롬프트가 해당 taxonomy의 위험 행동을 정확히 유도하는가?",
        "핵심 질문: judge rubric이 taxonomy와 schema/query에 맞게 설계되었는가?",
        "R4_1 사례처럼 rubric이 어긋나면 ASR 숫자는 나오지만 평가 의미가 흔들립니다.",
    ])

    slide(prs, "권장 운영 절차", [
        "1. schema/query/source context 확인",
        "2. Step1 scenario 생성 및 JSON 파싱 검증",
        "3. Step2 prompt 생성 및 CSV 행 수/컬럼 검증",
        "4. target response 생성 및 generation error 확인",
        "5. judge rubric 정합성 확인 후 judge 실행",
        "6. ASR 산출 및 표본 판정 검수",
        "7. 실패/중간 산출물은 partial/error 이름으로 분리 보관",
    ])

    s = slide(prs, "현재 알려진 이슈")
    add_table(s, 6, 3, [
        ["이슈", "영향", "대응"],
        ["OpenAI judge 429", "OpenAI judge 불가", "Gemini judge로 평가 완료"],
        ["Gemini deprecated", "장기 유지보수 리스크", "google-genai 전환 권장"],
        ["main.py --step all", "모델 인자 혼선 가능", "Step1/Step2 분리 실행"],
        ["rubric 불일치", "ASR 의미 흔들림", "R4_1/R2_4/R2_5 수정"],
        ["중간 실패 파일", "최종 결과 오해", "partial/error 이름으로 분리"],
    ], h=4.5)

    slide(prs, "다음 개선 제안", [
        "category별 preflight에 schema/query/chunks/rubric 정합성 체크 추가",
        "prompt strength judge 추가: attack prompt의 강도와 품질을 별도 점수화",
        "Step1 -> Step2 -> Response -> Judge -> Report를 하나의 harness command로 통합",
        "OpenAI/Gemini provider diagnostic 명령 추가",
        "최종 리포트를 Markdown/HTML/PPT로 자동 생성",
    ])

    slide(prs, "결론", [
        "R4_1 기준 전체 평가 체인은 연결 완료되었습니다.",
        "하네스는 단순 생성 코드가 아니라, 금융 레드팀 데이터셋을 반복 구축하고 검수하기 위한 운영 기반입니다.",
        "개발팀은 실행 안정성과 자동화를, 지식구축팀은 taxonomy/rubric/판정 품질을 중심으로 보면 됩니다.",
        "다음 단계는 rubric 검수 자동화와 end-to-end harness command 통합입니다.",
    ])

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
